"""
PPTX 解析器：使用 python-pptx 提取每页的标题、正文、列表、表格与演讲者备注。

特性:
- 结构提取：幻灯片标题 → HEADING；正文 → PARAGRAPH；列表项 → LIST；表格 → TABLE
- 列表检测：通过段落底层 XML 的 buChar/buAutoNum 判定项目符号，或依据段落缩进层级
- 演讲者备注：作为带标注的段落保留，便于 RAG 利用隐藏讲解信息
- 图片占位：以 [图片] 标记保留结构位置（不做 OCR，与 DOCX/PDF 解析器保持一致策略）
"""

import re

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn
from models.document import Document, ContentBlock, BlockType
from parsers.base import BaseParser


def _is_bulleted(para) -> bool:
    """通过底层 XML 判断段落是否带项目符号。

    注意：部分 python-pptx 版本（如 1.0.x）的 CT_TextParagraphProperties
    并未暴露 buChar/buAutoNum 等属性，直接访问会抛 AttributeError。
    因此这里直接查 <a:pPr> 下的子元素，兼容各版本。
    """
    try:
        pPr = para._p.find(qn("a:pPr"))
    except Exception:
        return False
    if pPr is None:
        return False
    if pPr.find(qn("a:buNone")) is not None:
        return False
    if pPr.find(qn("a:buChar")) is not None or pPr.find(qn("a:buAutoNum")) is not None:
        return True
    return False


class PptxParser(BaseParser):
    """使用 python-pptx 解析 PPTX 文件，提取结构、列表与表格信息"""

    def parse(self, file_path: str) -> Document:
        prs = Presentation(file_path)
        all_blocks: list[ContentBlock] = []

        for idx, slide in enumerate(prs.slides, start=1):
            slide_blocks: list[ContentBlock] = []

            # 幻灯片标题（顶层 title 占位符）
            title_shape = slide.shapes.title
            title_text = ""
            if title_shape is not None and title_shape.has_text_frame:
                title_text = title_shape.text_frame.text.strip()

            # 遍历所有形状
            for shape in slide.shapes:
                # 跳过标题形状（已在上方单独处理）
                if title_shape is not None and shape is title_shape:
                    continue

                # 表格
                if shape.has_table:
                    cb = self._table_to_block(shape.table)
                    if cb is not None:
                        slide_blocks.append(cb)
                    continue

                # 图片：以占位标记保留结构位置
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    slide_blocks.append(
                        ContentBlock(block_type=BlockType.PARAGRAPH, text="[图片]")
                    )
                    continue

                # 文本形状
                if not shape.has_text_frame:
                    continue
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if not text:
                        continue
                    is_list = _is_bulleted(para) or (para.level or 0) > 0
                    if is_list:
                        slide_blocks.append(
                            ContentBlock(block_type=BlockType.LIST, text=text)
                        )
                    else:
                        is_bold = any(r.font.bold for r in para.runs if r.font.bold)
                        is_italic = any(r.font.italic for r in para.runs if r.font.italic)
                        slide_blocks.append(
                            ContentBlock(
                                block_type=BlockType.PARAGRAPH,
                                text=text,
                                bold=is_bold,
                                italic=is_italic,
                            )
                        )

            # 演讲者备注
            notes_text = self._get_notes(slide)
            if notes_text:
                slide_blocks.append(
                    ContentBlock(
                        block_type=BlockType.PARAGRAPH,
                        text=f"[演讲者备注] {notes_text}",
                        italic=True,
                    )
                )

            # 跳过完全为空的幻灯片
            if not slide_blocks and not title_text:
                continue

            # 幻灯片标题（统一为 H1，便于检索时区分幻灯片边界）
            if title_text:
                all_blocks.append(
                    ContentBlock(
                        block_type=BlockType.HEADING,
                        text=f"Slide {idx}: {title_text}",
                        level=1,
                    )
                )
            elif slide_blocks:
                all_blocks.append(
                    ContentBlock(
                        block_type=BlockType.HEADING,
                        text=f"Slide {idx}",
                        level=1,
                    )
                )
            all_blocks.extend(slide_blocks)

        # 渲染为 Markdown 全文
        content = "\n\n".join(
            cb.to_markdown() for cb in all_blocks if not cb.is_empty
        )
        content = self._clean_text(content)

        return Document.from_file_path(file_path, content, content_blocks=all_blocks)

    # ── 表格转换 ──

    @staticmethod
    def _table_to_block(table) -> ContentBlock | None:
        """将 python-pptx 表格转为 ContentBlock"""
        all_rows: list[list[str]] = []
        for row in table.rows:
            cells = [c.text.strip().replace("\n", " ") for c in row.cells]
            if any(cells):
                all_rows.append(cells)
        if not all_rows:
            return None
        headers = all_rows[0]
        body = all_rows[1:] if len(all_rows) > 1 else []
        return ContentBlock(
            block_type=BlockType.TABLE,
            headers=list(headers),
            rows=body,
        )

    # ── 演讲者备注 ──

    @staticmethod
    def _get_notes(slide) -> str:
        """返回幻灯片的演讲者备注文本（无则返回空）"""
        try:
            notes_slide = slide.notes_slide
            text = notes_slide.notes_text_frame.text.strip()
            return text
        except Exception:
            return ""

    # ── 工具 ──

    @staticmethod
    def supported_extensions() -> list[str]:
        return [".pptx"]

    @staticmethod
    def _clean_text(text: str) -> str:
        """清理多余空白：合并连续换行，去除多余空格"""
        text = re.sub(r"\n{3,}", "\n\n", text)
        text = re.sub(r"[ \t]{2,}", " ", text)
        return text.strip()
