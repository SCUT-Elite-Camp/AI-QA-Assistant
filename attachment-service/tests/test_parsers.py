from __future__ import annotations

from pathlib import Path
import sys
from types import SimpleNamespace

from PIL import Image

from attachment_service import parsers


def test_csv_detects_non_utf8_encoding_and_exact_range(tmp_path: Path) -> None:
    path = tmp_path / "销售.csv"
    path.write_bytes("名称,金额\n华南,100\n".encode("gb18030"))
    items = parsers.parse_attachment(path, "att_csv", ".csv")
    assert "华南" in items[0]["content"]
    assert items[0]["locator"] == {"sheet": "CSV", "cell_range": "A1:B2"}


def test_html_extracts_only_static_visible_text_without_active_or_remote_content(tmp_path: Path) -> None:
    path = tmp_path / "notice.html"
    path.write_text(
        """<!doctype html><html><head>
        <style>.hidden{display:none}</style><script>fetch('https://outside.invalid/secret')</script>
        <link rel="stylesheet" href="https://outside.invalid/site.css"></head><body>
        <h1>内部通知</h1><img src="https://outside.invalid/tracker.png" alt="远程图片">
        <iframe src="https://outside.invalid/frame">外部框架</iframe>
        <object data="https://outside.invalid/object">对象内容</object>
        <embed src="https://outside.invalid/embed">正文内容</body></html>""",
        encoding="utf-8",
    )
    items = parsers.parse_attachment(path, "att_html", ".html")
    content = items[0]["content"]
    assert "内部通知" in content and "正文内容" in content
    assert "fetch(" not in content
    assert "outside.invalid" not in content
    assert "外部框架" not in content and "对象内容" not in content


def test_xlsx_preserves_sheet_range_and_marks_uncached_formula(tmp_path: Path) -> None:
    import openpyxl
    path = tmp_path / "report.xlsx"
    workbook = openpyxl.Workbook()
    sheet = workbook.active
    sheet.title = "汇总"
    sheet.append(["项目", "金额"])
    sheet.append(["总计", "=1+2"])
    workbook.save(path)
    items = parsers.parse_attachment(path, "att_xlsx", ".xlsx")
    assert items[0]["locator"] == {"sheet": "汇总", "cell_range": "A1:B2"}
    assert items[0]["confidence"] == 0.7


def test_docx_keeps_block_and_page_locators(tmp_path: Path) -> None:
    from docx import Document
    path = tmp_path / "policy.docx"
    document = Document()
    document.add_paragraph("第一页制度")
    document.add_page_break()
    document.add_paragraph("第二页流程")
    table = document.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "字段"
    table.cell(1, 0).text = "值"
    document.save(path)
    items = parsers.parse_attachment(path, "att_docx", ".docx")
    assert any(item["content"] == "第一页制度" and item["locator"]["page"] == 1 for item in items)
    assert any(item["content"] == "第二页流程" and item["locator"]["page"] == 2 for item in items)
    assert any(item["source_type"] == "table" for item in items)


def test_docx_embedded_image_ocr_inherits_body_page_locator(tmp_path: Path, monkeypatch) -> None:
    from docx import Document
    from docx.shared import Inches

    image_path = tmp_path / "error.png"
    Image.new("RGB", (30, 20), "white").save(image_path)
    path = tmp_path / "with-image.docx"
    document = Document()
    document.add_paragraph("第一页")
    document.add_page_break()
    document.add_picture(str(image_path), width=Inches(1))
    document.save(path)
    monkeypatch.setattr(parsers, "_ocr", lambda _image, attachment_id, locator: [
        parsers._item(attachment_id, "ocr_text", "DB-1042", "fake-ocr", locator, 0.95)
    ])

    items = parsers.parse_attachment(path, "att_docx_image", ".docx")
    embedded = next(item for item in items if item["content"] == "DB-1042")
    assert embedded["locator"]["page"] == 2
    assert embedded["locator"]["embedded_image"] == 1


def test_pptx_preserves_slide_and_ocr_locator(tmp_path: Path, monkeypatch) -> None:
    from pptx import Presentation
    from pptx.util import Inches
    image_path = tmp_path / "screen.png"
    Image.new("RGB", (20, 20), "white").save(image_path)
    path = tmp_path / "slides.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(1))
    box.text = "故障处理"
    slide.shapes.add_picture(str(image_path), Inches(1), Inches(2))
    presentation.save(path)

    monkeypatch.setattr(parsers, "_ocr", lambda _image, attachment_id, locator: [
        parsers._item(attachment_id, "ocr_text", "DB-1042", "fake-ocr", locator, 0.9)
    ])
    items = parsers.parse_attachment(path, "att_pptx", ".pptx")
    assert any(item["content"] == "故障处理" and item["locator"]["slide"] == 1 for item in items)
    assert any(item["content"] == "DB-1042" and item["locator"]["slide"] == 1 for item in items)


def test_low_text_pdf_enters_ocr_with_page_locator(tmp_path: Path, monkeypatch) -> None:
    import fitz
    path = tmp_path / "scan.pdf"
    pdf = fitz.open()
    page = pdf.new_page()
    page.insert_text((72, 72), "少")
    pdf.save(path)
    pdf.close()
    monkeypatch.setattr(parsers, "_ocr", lambda _image, attachment_id, locator: [
        parsers._item(attachment_id, "ocr_text", "扫描结果", "fake-ocr", locator, 0.95)
    ])
    items = parsers.parse_attachment(path, "att_pdf", ".pdf")
    assert any(item["content"] == "扫描结果" and item["locator"]["page"] == 1 for item in items)


def test_ppstructure_preserves_reading_order_and_markdown_table(monkeypatch) -> None:
    calls = {}

    class FakePipeline:
        def __init__(self, **options):
            calls["options"] = options

        def predict(self, image, **options):
            calls["predict"] = options
            return [SimpleNamespace(json={"res": {
                "parsing_res_list": [
                    {"block_id": 3, "block_order": 1, "block_label": "text", "block_content": "第一段", "block_bbox": [10, 20, 90, 40]},
                    {"block_id": 4, "block_order": 2, "block_label": "table", "block_content": "|字段|值|\n|---|---|", "block_bbox": [10, 50, 90, 90]},
                ],
                "overall_ocr_res": {"rec_texts": ["第一段"], "rec_scores": [0.98], "rec_boxes": [[10, 20, 90, 40]]},
            }})]

    monkeypatch.setitem(sys.modules, "paddleocr", SimpleNamespace(PPStructureV3=FakePipeline))
    monkeypatch.setenv("ATTACHMENT_OFFLINE_MODE", "false")
    monkeypatch.setattr(parsers, "_OCR_ENGINE", None)
    items = parsers._ocr(Image.new("RGB", (100, 100), "white"), "att_ocr", {"page": 1})
    table = next(item for item in items if item["source_type"] == "table")
    assert table["content"].startswith("|字段|值|")
    assert table["locator"] == {
        "page": 1, "bbox": [0.1, 0.5, 0.9, 0.9], "block": 4, "reading_order": 2,
    }
    assert calls["options"]["device"] == "cpu"
    assert calls["options"]["enable_mkldnn"] is False
    assert calls["predict"] == {"format_block_content": True}
