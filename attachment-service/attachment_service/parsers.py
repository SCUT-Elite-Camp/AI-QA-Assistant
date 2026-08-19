from __future__ import annotations

import csv
import json
import os
import sys
import subprocess
from pathlib import Path
from typing import Any
from uuid import uuid4

from PIL import Image, ImageOps

from .office import remove_office_tree, resolve_libreoffice


_OCR_ENGINE: Any = None


def _item(attachment_id: str, source_type: str, content: str, parser: str,
          locator: dict[str, Any] | None = None, confidence: float | None = None) -> dict[str, Any]:
    return {
        "evidence_id": f"aev_{uuid4().hex}", "attachment_id": attachment_id,
        "source_type": source_type, "content": content.strip(), "locator": locator or {},
        "confidence": confidence, "parser": parser,
    }


def parse_attachment(path: Path, attachment_id: str, extension: str) -> list[dict[str, Any]]:
    if extension in {".png", ".jpg", ".jpeg", ".webp", ".bmp", ".tif", ".tiff"}:
        return _parse_image(path, attachment_id)
    if extension == ".pdf":
        return _parse_pdf(path, attachment_id)
    if extension in {".csv", ".txt", ".md", ".json", ".html", ".htm"}:
        return _parse_textual(path, attachment_id, extension)
    if extension in {".xlsx", ".xls"}:
        return _parse_spreadsheet(path, attachment_id, extension)
    if extension == ".pptx":
        return _parse_pptx(path, attachment_id)
    if extension == ".docx":
        return _parse_docx(path, attachment_id)
    if extension in {".doc", ".ppt"}:
        converted = _convert_legacy_office(path, extension)
        try:
            return parse_attachment(converted, attachment_id, converted.suffix.lower())
        finally:
            if not remove_office_tree(converted.parent.parent):
                raise RuntimeError("legacy_office_cleanup_failed")
    return _parse_office(path, attachment_id)


def _ocr(image: Image.Image, attachment_id: str, locator: dict[str, Any]) -> list[dict[str, Any]]:
    global _OCR_ENGINE
    try:
        from paddleocr import PPStructureV3
    except ImportError as exc:
        raise RuntimeError("ocr_backend_unavailable") from exc
    config_path = os.getenv("PP_STRUCTURE_CONFIG_PATH", "").strip()
    offline = os.getenv("ATTACHMENT_OFFLINE_MODE", "true").lower() in {"1", "true", "yes"}
    if offline and (not config_path or not Path(config_path).is_file()):
        raise RuntimeError("ocr_backend_unavailable")
    options: dict[str, Any] = {
        "use_doc_orientation_classify": True,
        "use_doc_unwarping": True,
        "use_textline_orientation": True,
        "device": "cpu",
        # Paddle 3.3 on Windows currently fails in the oneDNN/PIR path for the
        # lightweight layout model. Keep the portable CPU executor by default.
        "enable_mkldnn": os.getenv(
            "PADDLE_OCR_ENABLE_MKLDNN", "false"
        ).lower() in {"1", "true", "yes"},
    }
    if config_path:
        options["paddlex_config"] = config_path
    if _OCR_ENGINE is None:
        _OCR_ENGINE = PPStructureV3(**options)
    import numpy as np
    results = _OCR_ENGINE.predict(np.asarray(image), format_block_content=True)
    items: list[dict[str, Any]] = []
    for result in results:
        payload = getattr(result, "json", None)
        if callable(payload):
            payload = payload()
        if not isinstance(payload, dict):
            continue
        data = payload.get("res", payload)
        structured_contents: set[str] = set()
        for block in data.get("parsing_res_list", []) if isinstance(data, dict) else []:
            if not isinstance(block, dict):
                continue
            content = str(block.get("block_content") or "").strip()
            if not content:
                continue
            label = str(block.get("block_label") or "text").lower()
            block_locator = dict(locator)
            normalized_box = _normalized_bbox(block.get("block_bbox"), *image.size)
            if normalized_box is not None:
                block_locator["bbox"] = normalized_box
            if block.get("block_id") is not None:
                block_locator["block"] = int(block["block_id"])
            if block.get("block_order") is not None:
                block_locator["reading_order"] = int(block["block_order"])
            source_type = "table" if "table" in label else "document_text"
            items.append(_item(
                attachment_id, source_type, content,
                "pp-structure-v3-lightweight", block_locator,
                _as_confidence(block.get("confidence") or block.get("score")),
            ))
            structured_contents.add(content)
        for ocr_data in _find_ocr_results(data):
            texts = ocr_data.get("rec_texts") or []
            scores = ocr_data.get("rec_scores") or []
            boxes = ocr_data.get("rec_boxes") or []
            width, height = image.size
            for index, text in enumerate(texts):
                if not str(text).strip():
                    continue
                box = boxes[index] if index < len(boxes) else None
                item_locator = dict(locator)
                normalized_box = _normalized_bbox(box, width, height)
                if normalized_box is not None:
                    item_locator["bbox"] = normalized_box
                score = float(scores[index]) if index < len(scores) else None
                items.append(_item(attachment_id, "ocr_text", str(text), "pp-structure-v3-lightweight", item_locator, score))
        for table in _find_table_results(data):
            content = str(
                table.get("pred_markdown")
                or table.get("markdown")
                or table.get("pred_html")
                or table.get("html")
                or ""
            ).strip()
            if content and content not in structured_contents:
                table_locator = dict(locator)
                box = table.get("bbox") or table.get("table_bbox")
                normalized_box = _normalized_bbox(box, *image.size)
                if normalized_box is not None:
                    table_locator["bbox"] = normalized_box
                items.append(_item(
                    attachment_id, "table", content,
                    "pp-structure-v3-lightweight", table_locator,
                    _as_confidence(table.get("confidence") or table.get("score")),
                ))
    return items


def _normalized_bbox(value: Any, width: int, height: int) -> list[float] | None:
    if value is None or width <= 0 or height <= 0:
        return None
    try:
        coordinates = list(value)
    except TypeError:
        return None
    if len(coordinates) < 4:
        return None
    try:
        if hasattr(coordinates[0], "__len__") and not isinstance(coordinates[0], (str, bytes)):
            points = [list(point) for point in coordinates]
            xs = [float(point[0]) for point in points if len(point) >= 2]
            ys = [float(point[1]) for point in points if len(point) >= 2]
            if not xs or not ys:
                return None
            x1, y1, x2, y2 = min(xs), min(ys), max(xs), max(ys)
        else:
            x1, y1, x2, y2 = [float(item) for item in coordinates[:4]]
    except (TypeError, ValueError):
        return None
    return [
        max(0.0, min(1.0, x1 / width)), max(0.0, min(1.0, y1 / height)),
        max(0.0, min(1.0, x2 / width)), max(0.0, min(1.0, y2 / height)),
    ]


def _find_ocr_results(value: Any):
    if isinstance(value, dict):
        if "rec_texts" in value:
            yield value
        for nested in value.values():
            yield from _find_ocr_results(nested)
    elif isinstance(value, list):
        for nested in value:
            yield from _find_ocr_results(nested)


def _find_table_results(value: Any):
    if isinstance(value, dict):
        if any(key in value for key in ("pred_markdown", "markdown", "pred_html", "html")) and (
            "table" in str(value.get("type", "table")).lower()
            or any(key in value for key in ("table_bbox", "pred_html"))
        ):
            yield value
        for nested in value.values():
            yield from _find_table_results(nested)
    elif isinstance(value, list):
        for nested in value:
            yield from _find_table_results(nested)


def _as_confidence(value: Any) -> float | None:
    try:
        return min(1.0, max(0.0, float(value)))
    except (TypeError, ValueError):
        return None


def _parse_image(path: Path, attachment_id: str) -> list[dict[str, Any]]:
    with Image.open(path) as image:
        normalized = ImageOps.exif_transpose(image).convert("RGB")
        return _ocr(normalized, attachment_id, {})


def _parse_pdf(path: Path, attachment_id: str) -> list[dict[str, Any]]:
    import fitz
    items: list[dict[str, Any]] = []
    with fitz.open(path) as pdf:
        for page_number, page in enumerate(pdf, 1):
            blocks = page.get_text("blocks")
            page_text_length = 0
            for block in blocks:
                text = str(block[4]).strip()
                if text:
                    page_text_length += len(text)
                    rect = page.rect
                    locator = {"page": page_number, "bbox": [block[0] / rect.width, block[1] / rect.height, block[2] / rect.width, block[3] / rect.height]}
                    items.append(_item(attachment_id, "document_text", text, "pymupdf", locator, 1.0))
            try:
                tables = page.find_tables()
                for table in tables.tables:
                    rows = table.extract()
                    content = "\n".join(" | ".join("" if cell is None else str(cell) for cell in row) for row in rows)
                    if content.strip():
                        rect = page.rect
                        box = table.bbox
                        items.append(_item(attachment_id, "table", content, "pymupdf-table", {
                            "page": page_number,
                            "bbox": [box[0] / rect.width, box[1] / rect.height, box[2] / rect.width, box[3] / rect.height],
                        }, 1.0))
            except (AttributeError, ValueError):
                pass
            if page_text_length < 30:
                pix = page.get_pixmap(matrix=fitz.Matrix(2, 2), alpha=False)
                image = Image.frombytes("RGB", (pix.width, pix.height), pix.samples)
                items.extend(_ocr(image, attachment_id, {"page": page_number}))
    return items


def _parse_textual(path: Path, attachment_id: str, extension: str) -> list[dict[str, Any]]:
    data = path.read_bytes()
    text = _decode_enterprise_text(data)
    if extension == ".html" or extension == ".htm":
        from bs4 import BeautifulSoup
        soup = BeautifulSoup(text, "html.parser")
        for node in soup(["script", "style", "iframe", "object", "embed"]):
            node.decompose()
        text = soup.get_text("\n")
    if extension == ".json":
        text = json.dumps(json.loads(text), ensure_ascii=False, indent=2)
    if extension == ".csv":
        rows = list(csv.reader(text.splitlines()))
        text = "\n".join(" | ".join(cell for cell in row) for row in rows)
        if not text.strip():
            return []
        from openpyxl.utils import get_column_letter
        columns = max((len(row) for row in rows), default=1)
        return [_item(attachment_id, "table", text, "python-csv", {
            "sheet": "CSV", "cell_range": f"A1:{get_column_letter(max(columns, 1))}{max(len(rows), 1)}"
        }, 1.0)]
    return [_item(attachment_id, "document_text", text, "builtin-text", {}, 1.0)] if text.strip() else []


def _decode_enterprise_text(data: bytes) -> str:
    if data.startswith((b"\xff\xfe", b"\xfe\xff")):
        return data.decode("utf-16")
    if data.startswith(b"\xef\xbb\xbf"):
        return data.decode("utf-8-sig")
    try:
        return data.decode("utf-8")
    except UnicodeDecodeError:
        pass
    # GB18030 is deterministic and covers GBK/GB2312, which remain common in
    # private Chinese enterprise exports. Only fall back to statistical
    # detection after this explicit decoding path.
    try:
        return data.decode("gb18030")
    except UnicodeDecodeError:
        from charset_normalizer import from_bytes
        match = from_bytes(data).best()
        if match is None:
            raise RuntimeError("text_encoding_undetected")
        return str(match)


def _parse_spreadsheet(path: Path, attachment_id: str, extension: str) -> list[dict[str, Any]]:
    items: list[dict[str, Any]] = []
    if extension == ".xlsx":
        import openpyxl
        values = openpyxl.load_workbook(path, read_only=True, data_only=True)
        formulas = openpyxl.load_workbook(path, read_only=True, data_only=False)
        try:
            for sheet in values.worksheets:
                formula_sheet = formulas[sheet.title]
                rows: list[list[str]] = []
                low_confidence = False
                max_column = 1
                for row_number, row in enumerate(sheet.iter_rows(values_only=True), 1):
                    cells = ["" if value is None else str(value) for value in row]
                    formula_values = next(formula_sheet.iter_rows(min_row=row_number, max_row=row_number, values_only=True))
                    if any(isinstance(value, str) and value.startswith("=") and not cells[index] for index, value in enumerate(formula_values)):
                        low_confidence = True
                    if any(cells):
                        max_column = max(max_column, len(cells))
                        rows.append(cells)
                if rows:
                    from openpyxl.utils import get_column_letter
                    content = "\n".join(" | ".join(row) for row in rows)
                    items.append(_item(attachment_id, "table", content, "openpyxl", {
                        "sheet": sheet.title, "cell_range": f"A1:{get_column_letter(max_column)}{sheet.max_row}"
                    }, 0.7 if low_confidence else 1.0))
        finally:
            values.close(); formulas.close()
        return items
    try:
        import xlrd
    except ImportError as exc:
        raise RuntimeError("xls_backend_unavailable") from exc
    workbook = xlrd.open_workbook(path, on_demand=True)
    try:
        for sheet in workbook.sheets():
            rows = [[str(sheet.cell_value(row, col)) for col in range(sheet.ncols)] for row in range(sheet.nrows)]
            content = "\n".join(" | ".join(row) for row in rows if any(row))
            if content:
                from openpyxl.utils import get_column_letter
                items.append(_item(attachment_id, "table", content, "xlrd", {
                    "sheet": sheet.name, "cell_range": f"A1:{get_column_letter(max(sheet.ncols, 1))}{max(sheet.nrows, 1)}"
                }, 1.0))
    finally:
        workbook.release_resources()
    return items


def _parse_pptx(path: Path, attachment_id: str) -> list[dict[str, Any]]:
    from io import BytesIO
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    presentation = Presentation(path)
    items: list[dict[str, Any]] = []
    for slide_number, slide in enumerate(presentation.slides, 1):
        texts: list[str] = []
        for shape in slide.shapes:
            if getattr(shape, "has_table", False):
                rows = [[cell.text.strip() for cell in row.cells] for row in shape.table.rows]
                content = "\n".join(" | ".join(row) for row in rows)
                if content.strip():
                    items.append(_item(attachment_id, "table", content, "python-pptx", {"slide": slide_number}, 1.0))
            elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                with Image.open(BytesIO(shape.image.blob)) as image:
                    items.extend(_ocr(image.convert("RGB"), attachment_id, {"slide": slide_number}))
            elif getattr(shape, "has_text_frame", False):
                text = shape.text.strip()
                if text:
                    texts.append(text)
        if texts:
            items.append(_item(attachment_id, "document_text", "\n".join(texts), "python-pptx", {"slide": slide_number}, 1.0))
    return items


def _parse_docx(path: Path, attachment_id: str) -> list[dict[str, Any]]:
    from io import BytesIO
    import hashlib
    from zipfile import ZipFile
    from docx import Document
    from docx.oxml.ns import qn
    from docx.table import Table
    from docx.text.paragraph import Paragraph

    document = Document(path)
    items: list[dict[str, Any]] = []
    page = 1
    block_index = 0
    embedded_image_index = 0
    positioned_image_hashes: set[str] = set()

    def append_embedded_images(element: Any, image_page: int) -> int:
        nonlocal embedded_image_index
        current_page = image_page
        for descendant in element.iter():
            if descendant.tag == qn("w:br") and descendant.get(qn("w:type")) == "page":
                current_page += 1
                continue
            if descendant.tag != qn("a:blip"):
                continue
            relationship_id = descendant.get(qn("r:embed"))
            related_part = document.part.related_parts.get(relationship_id) if relationship_id else None
            blob = getattr(related_part, "blob", None)
            if not blob:
                continue
            embedded_image_index += 1
            positioned_image_hashes.add(hashlib.sha256(blob).hexdigest())
            try:
                with Image.open(BytesIO(blob)) as image:
                    items.extend(_ocr(image.convert("RGB"), attachment_id, {
                        "page": current_page, "block": block_index,
                        "embedded_image": embedded_image_index,
                    }))
            except (OSError, ValueError):
                continue
        return current_page

    for child in document.element.body.iterchildren():
        if child.tag == qn("w:p"):
            paragraph = Paragraph(child, document)
            text = paragraph.text.strip()
            if text:
                items.append(_item(attachment_id, "document_text", text, "python-docx", {
                    "page": page, "block": block_index,
                }, 1.0))
            page = append_embedded_images(child, page)
        elif child.tag == qn("w:tbl"):
            table = Table(child, document)
            rows = [[cell.text.strip() for cell in row.cells] for row in table.rows]
            content = "\n".join(" | ".join(row) for row in rows)
            if content.strip():
                items.append(_item(attachment_id, "table", content, "python-docx", {
                    "page": page, "block": block_index,
                }, 1.0))
            page = append_embedded_images(child, page)
        block_index += 1

    # Header/footer and other package media may not have a body position. Keep
    # them available for OCR, but explicitly represent the unknown page rather
    # than inventing a locator.
    with ZipFile(path) as archive:
        media = sorted(name for name in archive.namelist() if name.startswith("word/media/") and not name.endswith("/"))
        for name in media:
            try:
                blob = archive.read(name)
                if hashlib.sha256(blob).hexdigest() in positioned_image_hashes:
                    continue
                embedded_image_index += 1
                with Image.open(BytesIO(blob)) as image:
                    items.extend(_ocr(image.convert("RGB"), attachment_id, {
                        "page": None, "embedded_image": embedded_image_index,
                    }))
            except (OSError, ValueError):
                continue
    return items


def _convert_legacy_office(path: Path, extension: str) -> Path:
    executable = resolve_libreoffice()
    if not executable:
        raise RuntimeError("legacy_office_backend_unavailable")
    output_dir = path.parent / f"office_convert_{uuid4().hex}"
    profile_dir = output_dir / "profile"
    converted_dir = output_dir / "output"
    profile_dir.mkdir(parents=True)
    converted_dir.mkdir()
    target_format = "docx" if extension == ".doc" else "pptx"
    result = subprocess.run(
        [
            executable,
            "--headless", "--nologo", "--nodefault", "--nofirststartwizard",
            "--nolockcheck", "--norestore",
            f"-env:UserInstallation={profile_dir.resolve().as_uri()}",
            "--convert-to", target_format, "--outdir", str(converted_dir), str(path),
        ],
        capture_output=True, timeout=120, check=False,
    )
    converted = converted_dir / f"{path.stem}.{target_format}"
    if result.returncode != 0 or not converted.exists():
        remove_office_tree(output_dir)
        raise RuntimeError("legacy_office_conversion_failed")
    return converted


def _parse_office(path: Path, attachment_id: str) -> list[dict[str, Any]]:
    root = Path(__file__).resolve().parents[2]
    pipeline = root / "data-pipeline"
    if str(pipeline) not in sys.path:
        sys.path.insert(0, str(pipeline))
    from parsers.registry import parse_file
    documents = parse_file(str(path))
    if not isinstance(documents, list):
        documents = [documents]
    items: list[dict[str, Any]] = []
    for document in documents:
        for index, block in enumerate(document.content_blocks):
            content = block.to_markdown()
            if content.strip():
                source_type = "table" if str(block.block_type) == "table" else "document_text"
                items.append(_item(attachment_id, source_type, content, "data-pipeline", {"block": index}, 1.0))
    return items
